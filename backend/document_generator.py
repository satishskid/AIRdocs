"""
Document Generation Service for GreyBrain Bank AI Aggregation Platform
Converts AI-generated content to professional formats: PDF, PowerPoint, Word
"""

import os
import io
import base64
from datetime import datetime
from typing import Dict, List, Optional, Union
import markdown
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from docx import Document
from docx.shared import Inches as DocxInches
from docx.enum.text import WD_ALIGN_PARAGRAPH
import re


class DocumentGenerator:
    """Professional document generation service"""
    
    def __init__(self):
        self.output_dir = "generated_documents"
        os.makedirs(self.output_dir, exist_ok=True)
        
        # GreyBrain Bank branding colors
        self.brand_colors = {
            'primary': RGBColor(255, 140, 0),  # Orange
            'secondary': RGBColor(64, 64, 64),  # Dark Gray
            'accent': RGBColor(245, 245, 245),  # Light Gray
            'text': RGBColor(33, 33, 33)  # Near Black
        }
    
    def generate_document(self, content: str, metadata: Dict, output_formats: List[str]) -> Dict:
        """Generate documents in specified formats"""
        results = {}
        
        # Parse content structure
        parsed_content = self._parse_content(content, metadata)
        
        # Generate each requested format
        for format_type in output_formats:
            try:
                if format_type.lower() == 'pdf':
                    file_path = self._generate_pdf(parsed_content, metadata)
                elif format_type.lower() in ['powerpoint', 'ppt', 'pptx']:
                    file_path = self._generate_powerpoint(parsed_content, metadata)
                elif format_type.lower() in ['word', 'doc', 'docx']:
                    file_path = self._generate_word(parsed_content, metadata)
                else:
                    continue
                    
                results[format_type] = {
                    'file_path': file_path,
                    'file_size': os.path.getsize(file_path),
                    'download_url': f"/download/{os.path.basename(file_path)}"
                }
            except Exception as e:
                results[format_type] = {'error': str(e)}
        
        return results
    
    def _parse_content(self, content: str, metadata: Dict) -> Dict:
        """Parse markdown content into structured format"""
        lines = content.split('\n')
        parsed = {
            'title': metadata.get('document_title', 'Generated Document'),
            'sections': [],
            'slides': [],
            'metadata': metadata
        }
        
        current_section = None
        current_slide = None
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            # Main title (# Title)
            if line.startswith('# ') and not line.startswith('## '):
                parsed['title'] = line[2:].strip()
                
            # Section headers (## Section)
            elif line.startswith('## '):
                if current_section:
                    parsed['sections'].append(current_section)
                current_section = {
                    'title': line[3:].strip(),
                    'content': [],
                    'subsections': []
                }
                
                # Also create slide for presentations
                if current_slide:
                    parsed['slides'].append(current_slide)
                current_slide = {
                    'title': line[3:].strip(),
                    'content': [],
                    'bullet_points': []
                }
                
            # Subsection headers (### Subsection)
            elif line.startswith('### '):
                if current_section:
                    subsection = {
                        'title': line[4:].strip(),
                        'content': []
                    }
                    current_section['subsections'].append(subsection)
                    
            # Bullet points
            elif line.startswith('• ') or line.startswith('- ') or line.startswith('* '):
                bullet_text = line[2:].strip()
                if current_section:
                    current_section['content'].append(f"• {bullet_text}")
                if current_slide:
                    current_slide['bullet_points'].append(bullet_text)
                    
            # Regular content
            elif line and not line.startswith('#'):
                if current_section:
                    current_section['content'].append(line)
                if current_slide:
                    current_slide['content'].append(line)
        
        # Add final section/slide
        if current_section:
            parsed['sections'].append(current_section)
        if current_slide:
            parsed['slides'].append(current_slide)
            
        return parsed
    
    def _generate_pdf(self, parsed_content: Dict, metadata: Dict) -> str:
        """Generate professional PDF document"""
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"{metadata.get('content_category', 'document')}_{timestamp}.pdf"
        file_path = os.path.join(self.output_dir, filename)
        
        doc = SimpleDocTemplate(file_path, pagesize=letter, topMargin=1*inch)
        styles = getSampleStyleSheet()
        
        # Custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Title'],
            fontSize=24,
            spaceAfter=30,
            textColor=colors.HexColor('#FF8C00'),
            alignment=TA_CENTER
        )
        
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading1'],
            fontSize=16,
            spaceAfter=12,
            textColor=colors.HexColor('#404040'),
            borderWidth=1,
            borderColor=colors.HexColor('#FF8C00'),
            borderPadding=5
        )
        
        story = []
        
        # Title page
        story.append(Paragraph(parsed_content['title'], title_style))
        story.append(Spacer(1, 20))
        
        # Metadata table
        metadata_data = [
            ['Category', metadata.get('content_category', 'N/A').title()],
            ['Quality Level', f"{metadata.get('quality_level', 'N/A')} (Premium)" if metadata.get('quality_level') == 3 else f"{metadata.get('quality_level', 'N/A')} (Balanced)"],
            ['AI Model', metadata.get('model_used', 'N/A')],
            ['Generated', datetime.now().strftime("%B %d, %Y")],
            ['Word Count', f"{metadata.get('word_count', 'N/A')} words"]
        ]
        
        metadata_table = Table(metadata_data, colWidths=[2*inch, 3*inch])
        metadata_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, -1), colors.HexColor('#F5F5F5')),
            ('TEXTCOLOR', (0, 0), (-1, -1), colors.black),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 0), (-1, -1), 10),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 12),
            ('GRID', (0, 0), (-1, -1), 1, colors.HexColor('#CCCCCC'))
        ]))
        
        story.append(metadata_table)
        story.append(PageBreak())
        
        # Content sections
        for section in parsed_content['sections']:
            story.append(Paragraph(section['title'], heading_style))
            story.append(Spacer(1, 12))
            
            for content_line in section['content']:
                if content_line.strip():
                    story.append(Paragraph(content_line, styles['Normal']))
                    story.append(Spacer(1, 6))
            
            # Subsections
            for subsection in section['subsections']:
                story.append(Paragraph(subsection['title'], styles['Heading2']))
                story.append(Spacer(1, 8))
                
                for content_line in subsection['content']:
                    if content_line.strip():
                        story.append(Paragraph(content_line, styles['Normal']))
                        story.append(Spacer(1, 6))
            
            story.append(Spacer(1, 20))
        
        # Footer
        story.append(Spacer(1, 30))
        footer_text = "Generated by GreyBrain Bank AI Aggregation Platform - Made by GreyBrain.ai"
        story.append(Paragraph(footer_text, styles['Normal']))
        
        doc.build(story)
        return file_path
    
    def _generate_powerpoint(self, parsed_content: Dict, metadata: Dict) -> str:
        """Generate professional PowerPoint presentation"""
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"{metadata.get('content_category', 'presentation')}_{timestamp}.pptx"
        file_path = os.path.join(self.output_dir, filename)
        
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = parsed_content['title']
        subtitle.text = f"Generated by GreyBrain Bank AI Platform\n{metadata.get('content_category', '').title()} • Quality Level {metadata.get('quality_level', 'N/A')}"
        
        # Style title slide
        title.text_frame.paragraphs[0].font.color.rgb = self.brand_colors['primary']
        subtitle.text_frame.paragraphs[0].font.color.rgb = self.brand_colors['secondary']
        
        # Content slides
        for slide_data in parsed_content['slides']:
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = slide_data['title']
            title.text_frame.paragraphs[0].font.color.rgb = self.brand_colors['primary']
            
            # Add bullet points
            text_frame = content.text_frame
            text_frame.clear()
            
            # Add content paragraphs
            for i, bullet in enumerate(slide_data['bullet_points'][:6]):  # Limit to 6 bullets per slide
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = bullet
                p.level = 0
                p.font.size = Pt(18)
                p.font.color.rgb = self.brand_colors['text']
            
            # Add additional content if no bullet points
            if not slide_data['bullet_points'] and slide_data['content']:
                content_text = '\n'.join(slide_data['content'][:3])  # Limit content
                if text_frame.paragraphs:
                    text_frame.paragraphs[0].text = content_text
                else:
                    p = text_frame.add_paragraph()
                    p.text = content_text
        
        # Final slide with branding
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add text box for branding
        left = Inches(2)
        top = Inches(3)
        width = Inches(6)
        height = Inches(2)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = "Thank You"
        
        p = text_frame.add_paragraph()
        p.text = "Generated by GreyBrain Bank AI Aggregation Platform"
        p.font.size = Pt(14)
        p.font.color.rgb = self.brand_colors['secondary']
        
        p = text_frame.add_paragraph()
        p.text = "Made by GreyBrain.ai"
        p.font.size = Pt(12)
        p.font.color.rgb = self.brand_colors['primary']
        
        prs.save(file_path)
        return file_path
    
    def _generate_word(self, parsed_content: Dict, metadata: Dict) -> str:
        """Generate professional Word document"""
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"{metadata.get('content_category', 'document')}_{timestamp}.docx"
        file_path = os.path.join(self.output_dir, filename)
        
        doc = Document()
        
        # Title
        title = doc.add_heading(parsed_content['title'], 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Metadata section
        doc.add_heading('Document Information', level=2)
        
        metadata_table = doc.add_table(rows=5, cols=2)
        metadata_table.style = 'Table Grid'
        
        metadata_items = [
            ('Category', metadata.get('content_category', 'N/A').title()),
            ('Quality Level', f"{metadata.get('quality_level', 'N/A')} ({'Premium' if metadata.get('quality_level') == 3 else 'Balanced'})"),
            ('AI Model', metadata.get('model_used', 'N/A')),
            ('Generated', datetime.now().strftime("%B %d, %Y")),
            ('Word Count', f"{metadata.get('word_count', 'N/A')} words")
        ]
        
        for i, (key, value) in enumerate(metadata_items):
            metadata_table.cell(i, 0).text = key
            metadata_table.cell(i, 1).text = str(value)
        
        doc.add_page_break()
        
        # Content sections
        for section in parsed_content['sections']:
            doc.add_heading(section['title'], level=1)
            
            for content_line in section['content']:
                if content_line.strip():
                    doc.add_paragraph(content_line)
            
            # Subsections
            for subsection in section['subsections']:
                doc.add_heading(subsection['title'], level=2)
                
                for content_line in subsection['content']:
                    if content_line.strip():
                        doc.add_paragraph(content_line)
        
        # Footer
        doc.add_paragraph()
        footer = doc.add_paragraph("Generated by GreyBrain Bank AI Aggregation Platform - Made by GreyBrain.ai")
        footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        doc.save(file_path)
        return file_path


# Global instance
document_generator = DocumentGenerator()
